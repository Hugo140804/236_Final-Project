# -*- coding: utf-8 -*-
"""Generate PPTX presentasi Blockchain API (16:9, tema navy-emas)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DIAGRAM = os.path.join(BASE, 'laporan', 'diagrams')
OUTPUT = os.path.join(BASE, 'laporan', 'Presentasi-Blockchain-API.pptx')

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x53, 0x95)
GOLD = RGBColor(0xE8, 0xA3, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
GRAY = RGBColor(0x59, 0x59, 0x59)
DARK = RGBColor(0x1C, 0x1E, 0x21)

SW, SH = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_rect(slide, left, top, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, left, top, w, h, text, size=18, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, left, top, w, h, items, size=17, color=DARK, gap=10, bullet='▪'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            head, rest = it
            r1 = p.add_run()
            r1.text = bullet + ' ' + head
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = BLUE
            r1.font.name = 'Calibri'
            r2 = p.add_run()
            r2.text = ' — ' + rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = 'Calibri'
        else:
            r = p.add_run()
            r.text = bullet + ' ' + it
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = 'Calibri'
    return tb


def header(slide, title, num):
    add_rect(slide, 0, 0, SW, 1.1, NAVY)
    add_rect(slide, 0, 1.1, SW, 0.06, GOLD)
    add_text(slide, 0.6, 0.18, 10.5, 0.8, title, size=28, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 11.9, 0.28, 1.0, 0.6, '%02d' % num, size=18, color=GOLD, bold=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, 7.28, SW, 0.22, LIGHT)
    add_text(slide, 0.5, 7.30, 8, 0.18, 'Blockchain API — Final Project 236', size=9, color=GRAY)


def png_size(path):
    with open(path, 'rb') as f:
        data = f.read(24)
    return int.from_bytes(data[16:20], 'big'), int.from_bytes(data[20:24], 'big')


def add_image_fit(slide, path, box_w, box_h, top_off=1.45):
    w, h = png_size(path)
    ratio = min(box_w / w, box_h / h)
    dw, dh = w * ratio, h * ratio
    left = (SW - dw) / 2
    top = top_off + (box_h - dh) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(dw), Inches(dh))


def add_table(slide, left, top, rows_data, col_widths, size=12):
    rows, cols = len(rows_data), len(rows_data[0])
    shp = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(sum(col_widths)), Inches(0.4 * rows))
    tbl = shp.table
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.name = 'Calibri'
                    r.font.color.rgb = WHITE if i == 0 else DARK
                    if i == 0:
                        r.font.bold = True
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
    return tbl

# ============ SLIDE 1: COVER ============
s = add_slide()
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 5.15, SW, 0.08, GOLD)
add_text(s, 0.8, 1.5, 11.7, 1.5, 'BLOCKCHAIN API', size=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 3.0, 11.7, 0.8, 'SaaS REST API — Data via API Key, Autentikasi JWT', size=24, color=GOLD,
         align=PP_ALIGN.CENTER)
add_text(s, 0.8, 5.5, 11.7, 1.2,
         'Express.js   •   PostgreSQL   •   Sequelize   •   Vercel\n\nFinal Project 236 — [NAMA] — [NIM]',
         size=18, color=LIGHT, align=PP_ALIGN.CENTER)

# ============ SLIDE 2: AGENDA ============
s = add_slide()
header(s, 'AGENDA', 2)
add_bullets(s, 1.2, 1.7, 10.9, 5, [
    'Latar Belakang & Rumusan Masalah',
    'Solusi & Fitur Aplikasi',
    'Teknologi yang Digunakan',
    'Use Case & User Flow',
    'ERD / Rancangan Database',
    'Endpoint API',
    'Demo Aplikasi (Postman)',
    'Hasil Pengujian',
    'Deployment di Vercel',
    'Kesimpulan',
], size=20, gap=11)

# ============ SLIDE 3: LATAR BELAKANG ============
s = add_slide()
header(s, 'LATAR BELAKANG & RUMUSAN MASALAH', 3)
add_text(s, 0.8, 1.5, 11.7, 0.5, 'Latar Belakang', size=20, color=BLUE, bold=True)
add_bullets(s, 1.2, 2.0, 11.3, 1.9, [
    'Data blockchain semakin dibutuhkan aplikasi, namun belum tersedia dalam bentuk API siap pakai.',
    'Layanan SaaS seperti OpenRouter / Weather API memberikan data lewat API key — pola inilah yang diadopsi.',
    'Membangun REST API yang menyediakan data kepada pihak lain dengan API key + JWT.',
], size=16, gap=8)
add_text(s, 0.8, 4.0, 11.7, 0.5, 'Rumusan Masalah', size=20, color=BLUE, bold=True)
add_bullets(s, 1.2, 4.5, 11.3, 2.2, [
    'Bagaimana membangun REST API SaaS berbasis API key?',
    'Bagaimana mengamankan akses data dengan JWT dan API key?',
    'Bagaimana menyediakan 50+ data kompleks dan men-deploy ke Vercel?',
], size=16, gap=8)

# ============ SLIDE 4: SOLUSI & FITUR ============
s = add_slide()
header(s, 'SOLUSI & FITUR', 4)
add_bullets(s, 1.2, 1.7, 11, 5.2, [
    ('Autentikasi JWT', 'login & manajemen akun pengembang (password di-hash bcrypt).'),
    ('Sistem API Key', 'konsumen cukup mengirim header x-api-key untuk mengambil data.'),
    ('CRUD Data', 'blockchain & kategori dengan relasi many-to-many.'),
    ('60 Data Blockchain', 'Bitcoin, Ethereum, Solana, Chainlink, dll. — melebihi minimal 50.'),
    ('13 Kategori', 'Layer 1, Layer 2, DeFi, NFT, Stablecoin, Meme, Oracle, dan lainnya.'),
    ('8 Pengembang', 'akun contoh untuk demo & pengujian.'),
], size=18, gap=13)

# ============ SLIDE 5: TEKNOLOGI ============
s = add_slide()
header(s, 'TEKNOLOGI YANG DIGUNAKAN', 5)
add_table(s, 1.0, 1.7, [
    ['Komponen', 'Teknologi', 'Fungsi'],
    ['Backend', 'Node.js + Express.js 5', 'Server REST API'],
    ['Database', 'PostgreSQL / Supabase', 'Penyimpanan data (5 tabel)'],
    ['ORM', 'Sequelize 6', 'Model, migration, seeder'],
    ['Autentikasi', 'JWT + bcrypt', 'Login & proteksi endpoint'],
    ['API Key', 'crypto.randomBytes', 'Kunci akses konsumen (blk_...)'],
    ['Deployment', 'Vercel', 'Hosting serverless'],
], [3.0, 4.0, 4.7], size=14)

# ============ SLIDE 6: USE CASE ============
s = add_slide()
header(s, 'USE CASE DIAGRAM', 6)
add_image_fit(s, os.path.join(DIAGRAM, 'usecase.png'), 11.6, 5.5)


# ============ SLIDE 7: USER FLOW ============
s = add_slide()
header(s, 'ACTIVITY DIAGRAM — USER FLOW', 7)
add_image_fit(s, os.path.join(DIAGRAM, 'activity-userflow.png'), 5.4, 5.5)
add_text(s, 8.2, 2.0, 4.5, 3.5,
         'Alur konsumen:\n\n1. Registrasi & login JWT\n2. Buat API key\n3. Kirim GET + x-api-key\n4. Key valid → data JSON\n5. Key tidak valid → 401',
         size=15, color=DARK)

# ============ SLIDE 8: ERD ============
s = add_slide()
header(s, 'ENTITY RELATIONSHIP DIAGRAM (ERD)', 8)
add_image_fit(s, os.path.join(DIAGRAM, 'erd.png'), 11.8, 5.6)

# ============ SLIDE 9: SKEMA DATABASE ============
s = add_slide()
header(s, '5 TABEL & RELASI', 9)
add_bullets(s, 1.0, 1.6, 11.3, 3.0, [
    ('pengembang', 'id, nama, email, password'),
    ('api_keys', 'id, pengembang_id, nama, key, aktif, terakhir_dipakai'),
    ('blockchain', 'id, nama, simbol, deskripsi, tahun_rilis, pengembang_id'),
    ('kategori', 'id, nama, deskripsi'),
    ('BlockchainKategori', 'blockchain_id + kategori_id (tabel relasi many-to-many)'),
], size=17, gap=9)
add_text(s, 1.0, 4.7, 11.3, 0.5, 'Relasi', size=20, color=BLUE, bold=True)
add_bullets(s, 1.2, 5.2, 11.3, 1.8, [
    'pengembang 1 — N blockchain (satu pengembang menulis banyak data)',
    'pengembang 1 — N api_keys (satu pengembang punya banyak key)',
    'blockchain N — N kategori (lewat tabel BlockchainKategori)',
], size=16, gap=8)

# ============ SLIDE 10: ENDPOINT ============
s = add_slide()
header(s, 'ENDPOINT API', 10)
add_table(s, 1.0, 1.6, [
    ['Method', 'Endpoint', 'Auth'],
    ['POST', '/api/register', '-'],
    ['POST', '/api/login', '-'],
    ['POST', '/api/apikey', 'JWT'],
    ['GET', '/api/blockchain', 'API Key / JWT'],
    ['POST', '/api/blockchain', 'JWT'],
    ['PUT', '/api/blockchain/:id', 'JWT'],
    ['DELETE', '/api/blockchain/:id', 'JWT'],
    ['GET', '/api/kategori', 'API Key / JWT'],
    ['POST / PUT / DELETE', '/api/kategori...', 'JWT'],
], [3.6, 6.0, 2.4], size=14)

# ============ SLIDE 11: DEMO ============
s = add_slide()
header(s, 'DEMO APLIKASI (POSTMAN)', 11)
add_bullets(s, 1.0, 1.6, 11.3, 5.3, [
    ('1. Register', 'POST /api/register → 201'),
    ('2. Login', 'POST /api/login → token JWT (auto-simpan)'),
    ('3. Cek akun', 'GET /api/me → id, nama, email'),
    ('4. Buat API key', 'POST /api/apikey → blk_xxxxxxxx'),
    ('5. Tambah data', 'POST /api/blockchain (nama, simbol, pengembang_id, kategori_ids)'),
    ('6. Akses konsumen', 'GET /api/blockchain + header x-api-key → 200, data JSON'),
], size=17, gap=11)

# ============ SLIDE 12: PENGUJIAN ============
s = add_slide()
header(s, 'HASIL PENGUJIAN', 12)
add_table(s, 1.0, 1.7, [
    ['Skenario', 'Hasil'],
    ['Register pengembang', '201 Created'],
    ['Login (email & password benar)', '200 + token JWT'],
    ['Akses data tanpa auth', '401 Unauthorized'],
    ['Akses data dengan API key valid', '200 + data'],
    ['Akses data dengan API key salah', '401'],
    ['POST / PUT / DELETE dengan JWT', '200 / 201'],
], [6.5, 5.2], size=14)

# ============ SLIDE 13: DEPLOYMENT ============
s = add_slide()
header(s, 'DEPLOYMENT DI VERCEL', 13)
add_bullets(s, 1.0, 1.7, 11.3, 4.2, [
    ('URL publik', 'https://236-final-project.vercel.app'),
    ('Cara deploy', 'hubungkan repo GitHub ke Vercel → otomatis build & deploy.'),
    ('Environment variables', 'POSTGRES_URL, JWT_SECRET, JWT_EXPIRES.'),
    ('Serverless', 'tanpa kelola server, auto-redeploy tiap push.'),
], size=18, gap=12)
add_text(s, 1.0, 5.6, 11.3, 0.8, 'Demo langsung di Postman dengan base URL Vercel di atas.', size=16,
         color=GOLD, bold=True)

# ============ SLIDE 14: KESIMPULAN ============
s = add_slide()
header(s, 'KESIMPULAN', 14)
add_bullets(s, 1.0, 1.7, 11.3, 5.0, [
    ('SaaS API berhasil', 'data disediakan ke konsumen lewat API key (register → login JWT → buat key → akses).'),
    ('Keamanan ganda', 'JWT untuk pengelolaan, API key untuk konsumsi data.'),
    ('60 data & 5 tabel', 'relasi many-to-many — melebihi syarat minimal 50 data & 2 tabel.'),
    ('Deploy di Vercel', 'dapat diakses publik di https://236-final-project.vercel.app'),
], size=18, gap=14)

# ============ SLIDE 15: TERIMA KASIH ============
s = add_slide()
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, 4.4, SW, 0.08, GOLD)
add_text(s, 0.8, 2.6, 11.7, 1.2, 'TERIMA KASIH', size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 4.0, 11.7, 0.7, 'Sesi Tanya Jawab', size=24, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 5.6, 11.7, 0.7, 'Blockchain API — Final Project 236', size=16, color=LIGHT,
         align=PP_ALIGN.CENTER)

# ============ SIMPAN ============
prs.core_properties.title = 'Presentasi Blockchain API - Final Project 236'
prs.core_properties.author = '236 Final Project'
prs.save(OUTPUT)
print('PPTX berhasil dibuat:', OUTPUT)
print('Jumlah slide:', len(prs.slides._sldIdLst))

